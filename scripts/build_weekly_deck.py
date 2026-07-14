"""Build the weekly progress deck (7 slides, CMU red/white, English).

Covers: variables (X/Y/controls), methodology, results (survival ladder,
pricing/margins, heterogeneity), the narrative, and status. All numbers real.
Output: reports/Weekly_Update.pptx

Usage:
    python scripts/build_weekly_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

RED = RGBColor(0xC4, 0x12, 0x30)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT = RGBColor(0xF4, 0xF1, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ASSETS = Path("reports/deck_assets")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = prs.slide_width


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def set_text(tf, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"


def title_bar(s, text, kicker=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background(); bar.shadow.inherit = False
    tb = box(s, 0.55, 0.1, 12.2, 0.95); tf = tb.text_frame; tf.word_wrap = True
    if kicker:
        p = tf.paragraphs[0]; r = p.add_run(); r.text = kicker.upper()
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RGBColor(0xF7, 0xC9, 0xD0); r.font.name = "Calibri"
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    r = p2.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"


def footer(s, n):
    set_text(box(s, 0.55, 7.06, 9, 0.35).text_frame,
             "CMU Capstone × Larridin   |   Weekly Update — July 2026", 9, GREY)
    set_text(box(s, 12.4, 7.06, 0.6, 0.35).text_frame, str(n), 9, GREY, align=PP_ALIGN.RIGHT)


def bullets(s, items, l, t, w, h, size=14, gap=7):
    tb = box(s, l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = it[0] if isinstance(it, tuple) else 0
        txt = it[1] if isinstance(it, tuple) else it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        r.font.name = "Calibri"


def set_cell(cell, text, size, color, bold, fill, align):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"


def table(s, data, l, t, col_w, fsize=12, rowh=0.42):
    nr, nc = len(data), len(data[0])
    gf = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(sum(col_w)), Inches(rowh * nr))
    tbl = gf.table; tbl.first_row = False; tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            is_h = i == 0
            fill = RED if is_h else (LIGHT if i % 2 == 0 else WHITE)
            set_cell(tbl.cell(i, j), str(val), fsize, WHITE if is_h else DARK, is_h, fill,
                     PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    return gf


def chip(s, l, t, w, h, text, fill, txtcolor=WHITE, size=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, text, size, txtcolor, bold=True, align=PP_ALIGN.CENTER)
    return sh


def pic(s, name, l, t, w=None):
    kw = {"width": Inches(w)} if w else {}
    return s.shapes.add_picture(str(ASSETS / name), Inches(l), Inches(t), **kw)


# ------------------------------------------------ Slide 1 — title/status
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), SW, Inches(2.2))
band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background(); band.shadow.inherit = False
set_text(box(s, 0.8, 2.5, 11.7, 1.1).text_frame, "Weekly Update: Analysis Complete, Paper Drafted", 34, WHITE, bold=True)
set_text(box(s, 0.8, 3.6, 11.7, 0.6).text_frame, "Do AI Adoption Signals Predict Company Performance?", 18, WHITE)
chip(s, 0.8, 5.0, 3.6, 0.55, "Data collection  DONE", LIGHT, DARK, 13)
chip(s, 4.6, 5.0, 3.6, 0.55, "Core analysis  DONE", LIGHT, DARK, 13)
chip(s, 8.4, 5.0, 3.9, 0.55, "Paper  FULL DRAFT (12 pp)", LIGHT, DARK, 13)

# ------------------------------------------------ Slide 2 — variables
s = slide(); title_bar(s, "The Variables", "Design")
set_text(box(s, 0.55, 1.25, 12, 0.4).text_frame, "X — AI adoption signals (4 families)", 15, RED, bold=True)
table(s, [
    ["Signal", "Source", "Coverage", "What it measures"],
    ["Larridin scores (adoption / proficiency / impact / maturity)", "Larridin Tracker, Jan 2026", "514", "Composite evidence-based AI adoption, 1–5"],
    ["Narrative concreteness  (ours)", "10-K filings via LLM", "462", "Deployed, quantified AI vs. aspirational language"],
    ["Investment intensity  (ours)", "10-K filings via LLM", "447", "Scale of AI commitment — building or buying"],
    ["AI-hiring builder rate  (ours)", "31k job postings, Jun+Jul", "253", "Share of hiring that builds AI systems"],
], 0.55, 1.7, [4.1, 2.5, 1.0, 4.6], fsize=11.5, rowh=0.42)
set_text(box(s, 0.55, 3.95, 12, 0.4).text_frame, "Y — outcomes    &    Controls", 15, RED, bold=True)
table(s, [
    ["Outcome", "n", "Controls (every regression)"],
    ["Revenue growth YoY, Q1-2026  (primary)", "438", "Sector fixed effects (12)"],
    ["Operating-margin change YoY", "348", "Firm size (log market cap at score date)"],
    ["Forward stock returns, 1–4 months", "506", "Growth momentum (pre-signal revenue growth)"],
], 0.55, 4.4, [4.8, 0.8, 6.6], fsize=11.5, rowh=0.42)
set_text(box(s, 0.55, 6.25, 12.2, 0.6).text_frame,
         "Design: cross-section — signals measured Jan–Mar 2026, outcomes realized afterward. Signals enter as percentile ranks.", 12, GREY)
footer(s, 2)

# ------------------------------------------------ Slide 3 — methodology
s = slide(); title_bar(s, "Methodology: Make the Signal Prove Itself", "Design")
bullets(s, [
    "Specification ladder — each signal must survive four increasingly demanding tests:",
    (1, "(1) raw association  →  (2) + sector effects  →  (3) + firm size  →  (4) + growth momentum"),
    "The momentum control is the key test: growing firms adopt more AI AND keep growing — spec (4) separates the two",
    "Rigor stack:",
    (1, "HC3 robust standard errors; outcomes winsorized 1%/99%; signals rank-transformed"),
    (1, "Multiple-testing control (Benjamini–Hochberg FDR) across all signal × outcome tests"),
    (1, "Classifier validated vs. 657 independent labels (~90%); 87–90% of LLM evidence quotes verify verbatim"),
    (1, "Complete-case robustness + per-cell sample sizes reported everywhere"),
], 0.55, 1.4, 12.2, 5.2, size=15, gap=11)
footer(s, 3)

# ------------------------------------------------ Slide 4 — headline result
s = slide(); title_bar(s, "Headline: Concreteness Survives Everything", "Results · 1 of 3")
pic(s, "survival_ladder.png", 0.55, 1.45, w=7.3)
bullets(s, [
    "All signals associate with revenue growth raw — the phenomenon is real",
    "Composite scores attenuate once size enters (they aggregate scale-correlated parts)",
    "Narrative concreteness survives all controls:",
    (1, "+8.7pp revenue-growth gap, low → high  (p = 0.007)"),
    (1, "passes FDR in the pre-specified primary family"),
    "Sorts are monotone: conc 2 → 5 bins = 7.3% → 52% growth",
], 8.1, 1.5, 4.7, 5.2, size=13, gap=9)
footer(s, 4)

# ------------------------------------------------ Slide 5 — where it shows / doesn't
s = slide(); title_bar(s, "Growth Channel, Not (Yet) Margins or Prices", "Results · 2 of 3")
table(s, [
    ["Outcome", "Result", "Interpretation"],
    ["Revenue growth", "Significant (conc +8.7pp***)", "AI adoption is visible in top-line fundamentals"],
    ["Operating margins", "Null across all signals", "AI is a growth story, not (yet) a cost story"],
    ["Stock returns (4-mo)", "Null after FDR", "Public signals already priced — consistent with market efficiency"],
], 0.55, 1.6, [2.6, 3.6, 6.0], fsize=13, rowh=0.55)
set_text(box(s, 0.55, 4.1, 12, 0.4).text_frame, "Value-chain heterogeneity (S&P 500, 4-category classification)", 15, RED, bold=True)
bullets(s, [
    "AI-Infrastructure suppliers: +48.9% mean 4-month return; +36.8pp vs. peers with sector & size controls (p<0.0001)",
    "Concreteness→growth link significant WITHIN Physical-Asset/Late Adopters (p=0.03, n=214) —",
    (1, "the signal works best precisely where AI claims are cheapest to make and hardest to verify"),
], 0.55, 4.55, 12.2, 2.2, size=14, gap=9)
footer(s, 5)

# ------------------------------------------------ Slide 6 — the story
s = slide(); title_bar(s, "The Story We Can Tell", "Narrative")
for i, (head, body) in enumerate([
    ("1. The thesis holds", "Evidence-based AI adoption measurement works: every signal family associates with real revenue growth. Top-quintile firms grew ~2x bottom-quintile."),
    ("2. Depth is the edge", "Surface AI mentions carry nothing; composite scores capture the phenomenon; the deepest evidence-anchored dimension — disclosure concreteness — survives every control. Structured measurement is exactly what separates signal from noise."),
    ("3. The product roadmap", "The concreteness dimension plugs straight into the Tracker; the hiring pipeline scales the POC to 536 companies monthly; each new score vintage compounds the dataset toward causal designs."),
]):
    chip(s, 0.55, 1.5 + i * 1.75, 3.0, 1.4, head, RED, size=14)
    tb = box(s, 3.8, 1.5 + i * 1.75, 8.9, 1.5); tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, body, 13, DARK)
set_text(box(s, 0.55, 6.55, 12.2, 0.5).text_frame,
         "One-liner: “CMU-validated — evidence-based AI adoption measurement predicts real revenue growth, and measurement depth is what separates signal from noise.”",
         12, GREY)
footer(s, 6)

# ------------------------------------------------ Slide 7 — status & next
s = slide(); title_bar(s, "Status & Next Steps", "Plan")
cw = 5.8
chip(s, 0.55, 1.45, cw, 0.5, "DONE THIS WEEK", RED, size=13)
bullets(s, [
    "July hiring snapshot: 536 companies, 17.3k postings (month 2 of panel)",
    "Signal reliability estimate: ρ = 0.54 month-over-month",
    "Controls pulled (size, momentum) — full regression suite run",
    "Value-chain segmentation integrated",
    "Full paper draft: 12 pages, 6 tables, compiled",
], 0.55, 2.1, cw, 4.2, size=13, gap=9)
chip(s, 6.9, 1.45, cw, 0.5, "NEXT", RGBColor(0x8A, 0x8A, 0x8A), size=13)
bullets(s, [
    "Team review of paper draft → revisions",
    "Final presentation deck",
    "Optional: productivity outcome (employee extraction)",
    "Stretch: dashboard wiring (deprioritized vs. paper)",
], 6.9, 2.1, cw, 4.2, size=13, gap=9)
footer(s, 7)

out = Path("reports/Weekly_Update.pptx")
prs.save(str(out))
print("saved ->", out, "| slides:", len(prs.slides._sldIdLst))
